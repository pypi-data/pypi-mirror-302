from __future__ import annotations
import json
from datetime import datetime, date, time, timedelta
from decimal import Decimal
from pathlib import Path
from dataclasses import dataclass
from PyPDF2 import PdfReader
from docx import Document as init_docx
from docx.document import Document as DocxDocument
from fitz import Document as PDFDocument
from fitz import open as open_pdf
from openpyxl import Workbook
from openpyxl.reader.excel import load_workbook as init_xlsx
from pptx import Presentation as init_pptx
from pptx.presentation import Presentation as PPTXDocument

from ._proxy import DocumentFile


class JsonEncoder(json.JSONEncoder):
    def default(self, o: object) -> object:
        if isinstance(o, datetime):
            return o.isoformat()
        if isinstance(o, date):
            return o.isoformat()
        if isinstance(o, time):
            return o.isoformat()
        if isinstance(o, timedelta):
            return o.total_seconds()
        if isinstance(o, Decimal):
            return str(o)
        return super().default(o)

@dataclass
class DocxFile(DocumentFile[DocxDocument]):
    def __load__(self):
        return init_docx(self.name)

    def extract_text(self):
        doc = self.__load__()
        for paragraph in doc.paragraphs:
            if paragraph.text:
                yield paragraph.text
            else:
                continue

    def extract_images(self):
        doc = self.__load__()
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    continue
                else:
                    for inline in run.element.iter():  # type: ignore
                        if inline.tag.endswith("inline"):  # type: ignore
                            for pic in inline.iter():  # type: ignore
                                if pic.tag.endswith("blip"):  # type: ignore
                                    image = pic.embed  # type: ignore
                                    image_part = run.part.related_parts[image]
                                    yield image_part.blob
                                else:
                                    continue
                        else:
                            continue


@dataclass
class PDFFile(DocumentFile[PDFDocument]):
    def __load__(self):
        return open_pdf(self.name)

    def extract_text(self):  # type: ignore
        text_doc = PdfReader(self.name)
        for page_number in range(len(text_doc.pages)):
            page = text_doc.pages[page_number]
            yield page.extract_text()

    def extract_images(self):
        img_doc = open_pdf(Path(self.name).as_posix())  # type: ignore
        for page in img_doc:  # type: ignore
            for img in page.get_images():  # type: ignore
                xref = img[0]  # type: ignore
                base_image = img_doc.extract_image(xref)  # type: ignore
                image_bytes = base_image["image"]  # type: ignore
                assert isinstance(image_bytes, bytes)
                yield image_bytes


@dataclass
class ExcelFile(DocumentFile[Workbook]):
    def __load__(self):
        return init_xlsx(self.name)

    def extract_text(self):
        wb = self.__load__()
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        data_dict = {
                            "row": cell.row,
                            "column": cell.column,
                            "value": cell.value,
                            "sheet": sheet_name,
                        }
                        yield json.dumps(data_dict, cls=JsonEncoder)

    def extract_images(self):
        wb = self.__load__()
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for image in sheet._images:  # type: ignore
                yield image._data()  # type: ignore


@dataclass
class PPTXFile(DocumentFile[PPTXDocument]):
    def __load__(self):
        return init_pptx(self.name)
    def extract_text(self):
        prs = self.__load__()
        for slide in prs.slides:  # type: ignore
            for shape in slide.shapes:  # type: ignore
                if shape.has_text_frame:  # type: ignore
                    text_frame = shape.text_frame  # type: ignore
                    for paragraph in text_frame.paragraphs:  # type: ignore
                        if paragraph.text:  # type: ignore
                            yield paragraph.text
                        else:
                            continue

    def extract_images(self):
        prs = self.__load__()
        for slide in prs.slides:  # type: ignore
            for shape in slide.shapes:  # type: ignore
                if shape.shape_type == 13:  # type: ignore
                    image = shape.image  # type: ignore
                    yield image.blob  # type: ignore
                else:
                    continue